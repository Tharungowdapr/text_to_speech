import os
import re
import tempfile

SUPPORTED = {".docx": "word", ".pptx": "powerpoint"}


class DocumentProcessor:
    @staticmethod
    def supported_formats():
        return list(SUPPORTED.keys())

    @staticmethod
    def is_supported(filename: str) -> bool:
        ext = os.path.splitext(filename)[1].lower()
        return ext in SUPPORTED

    @staticmethod
    def extract(filepath: str) -> str:
        ext = os.path.splitext(filepath)[1].lower()
        if ext == ".docx":
            return DocumentProcessor._extract_docx(filepath)
        elif ext == ".pptx":
            return DocumentProcessor._extract_pptx(filepath)
        return ""

    @staticmethod
    def extract_from_bytes(file_bytes: bytes, ext: str) -> str:
        """Extract text from document bytes for cross-instance support"""
        if ext == ".docx":
            return DocumentProcessor._extract_docx_from_bytes(file_bytes)
        elif ext == ".pptx":
            return DocumentProcessor._extract_pptx_from_bytes(file_bytes)
        return ""

    @staticmethod
    def _extract_docx(filepath: str) -> str:
        try:
            from docx import Document
            doc = Document(filepath)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            return f"[Error reading docx: {e}]"

    @staticmethod
    def _extract_docx_from_bytes(file_bytes: bytes) -> str:
        try:
            from docx import Document
            from io import BytesIO
            doc = Document(BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            return f"[Error reading docx: {e}]"

    @staticmethod
    def _extract_pptx(filepath: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:
            return f"[Error reading pptx: {e}]"

    @staticmethod
    def _extract_pptx_from_bytes(file_bytes: bytes) -> str:
        try:
            from pptx import Presentation
            from io import BytesIO
            prs = Presentation(BytesIO(file_bytes))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:
            return f"[Error reading pptx: {e}]"
